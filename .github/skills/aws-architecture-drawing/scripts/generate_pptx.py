#!/usr/bin/env python3
"""
AWS Architecture Diagram Generator (python-pptx)
Generates professional AWS architecture diagrams following the official
AWS Architecture Center guidelines (Release 2026.01.30).

This script implements the mandatory 8-step construction workflow:
  Step 1: Create Presentation & Blank Slide
  Step 2: Add Diagram Title
  Step 3: Add Groups (Outermost → Innermost)
  Step 4: Add Service Icons (0.83×0.83in)
  Step 5: Add Resource Icons (0.50×0.50in)
  Step 6: Add Arrows & Connectors
  Step 7: Add Annotations (Optional)
  Step 8: Save Output

Usage:
    # Import and use programmatically:
    from generate_pptx import AWSArchitectureDiagram
    diagram = AWSArchitectureDiagram("My Architecture", icon_base_path="path/to/icons")
    diagram.add_group("AWS Cloud", 1.0, 1.2, 11.0, 5.5)
    diagram.add_group("VPC", 1.5, 1.8, 10.0, 4.5, label="VPC (10.0.0.0/16)")
    diagram.add_service_icon("path/to/Arch_AWS-Lambda_48.png", "AWS Lambda", 5.0, 3.0)
    diagram.add_arrow(3.0, 3.5, 5.0, 3.5)
    diagram.save("output.pptx")
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn


# ═══════════════════════════════════════════════════════════════════════════════
# GROUP SPECIFICATIONS (from Slide 8-9 System Elements + Slide 25 Groups)
# ═══════════════════════════════════════════════════════════════════════════════

GROUP_SPECS = {
    "AWS Cloud": {
        "color": RGBColor(0x23, 0x2F, 0x3E),
        "dash": None,
        "icon": "AWS-Cloud_32.png",
    },
    "AWS Account": {
        "color": RGBColor(0xE7, 0x15, 0x7B),
        "dash": None,
        "icon": "AWS-Account_32.png",
    },
    "Region": {
        "color": RGBColor(0x00, 0xA4, 0xA6),
        "dash": MSO_LINE_DASH_STYLE.SQUARE_DOT,
        "icon": "Region_32.png",
    },
    "Availability Zone": {
        "color": RGBColor(0x00, 0xA4, 0xA6),
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "icon": None,
    },
    "VPC": {
        "color": RGBColor(0x8C, 0x4F, 0xFF),
        "dash": None,
        "icon": "Virtual-private-cloud-VPC_32.png",
    },
    "Public Subnet": {
        "color": RGBColor(0x7A, 0xA1, 0x16),
        "dash": None,
        "icon": "Public-subnet_32.png",
    },
    "Private Subnet": {
        "color": RGBColor(0x00, 0xA4, 0xA6),
        "dash": None,
        "icon": "Private-subnet_32.png",
    },
    "Security Group": {
        "color": RGBColor(0xDD, 0x34, 0x4C),
        "dash": None,
        "icon": None,
    },
    "Auto Scaling": {
        "color": RGBColor(0xED, 0x71, 0x00),
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "icon": "Auto-Scaling-group_32.png",
    },
    "EC2 Instance Contents": {
        "color": RGBColor(0xED, 0x71, 0x00),
        "dash": None,
        "icon": "EC2-instance-contents_32.png",
    },
    "Spot Fleet": {
        "color": RGBColor(0xED, 0x71, 0x00),
        "dash": None,
        "icon": "Spot-Fleet_32.png",
    },
    "Corporate Data Center": {
        "color": RGBColor(0x7D, 0x89, 0x98),
        "dash": None,
        "icon": "Corporate-data-center_32.png",
    },
    "Server Contents": {
        "color": RGBColor(0x7D, 0x89, 0x98),
        "dash": None,
        "icon": "Server-contents_32.png",
    },
    "IoT Greengrass": {
        "color": RGBColor(0x7A, 0xA1, 0x16),
        "dash": None,
        "icon": "AWS-IoT-Greengrass-Deployment_32.png",
    },
    "Step Functions": {
        "color": RGBColor(0xE7, 0x15, 0x7B),
        "dash": None,
        "icon": None,
    },
    "Generic": {
        "color": RGBColor(0x7D, 0x89, 0x98),  # Gray (official)
        "dash": None,
        "icon": None,
    },
}

# Icon base paths (relative to skill root: .github/skills/aws-architecture-drawing/)
SKILL_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS_DIR = os.path.join(SKILL_ROOT, "assets")
SERVICE_ICON_DIR = os.path.join(ASSETS_DIR, "service-icons")
RESOURCE_ICON_DIR = os.path.join(ASSETS_DIR, "resource-icons")
GROUP_ICON_DIR = os.path.join(ASSETS_DIR, "group-icons")
CATEGORY_ICON_DIR = os.path.join(ASSETS_DIR, "category-icons")
GENERAL_ICON_DIR = os.path.join(ASSETS_DIR, "general-icons")


class AWSArchitectureDiagram:
    """Professional AWS Architecture Diagram generator."""

    # Default icon sizes (inches)
    SERVICE_ICON_SIZE = 0.83
    RESOURCE_ICON_SIZE = 0.50
    GROUP_ICON_CLEARANCE = 0.55  # clearance from group corner icon
    # Label dimensions (proportional to icon size)
    SERVICE_LABEL_W = 1.30   # label textbox width — wider for 12pt text
    RESOURCE_LABEL_W = 1.00  # label textbox width for resource/general icons

    def __init__(self, title="AWS Architecture", icon_base_path=None, scale=1.0):
        """
        Initialize diagram.

        Args:
            title: Diagram title text
            icon_base_path: Override base path to assets directory.
            scale: Scale factor for all icons (0.5=half, 1.0=normal).
                   Use < 1.0 when diagram has many icons to prevent overflow.
        """
        self.assets_dir = icon_base_path or ASSETS_DIR
        self.group_icon_dir = os.path.join(self.assets_dir, "group-icons")
        self.service_icon_dir = os.path.join(self.assets_dir, "service-icons")
        self.resource_icon_dir = os.path.join(self.assets_dir, "resource-icons")
        self.general_icon_dir = os.path.join(self.assets_dir, "general-icons")
        self.scale = scale

        # Icon registry: name -> {cx, cy, w, h, type}
        # Used by connect() to anchor arrows to icon edges
        self._icons = {}

        # Arrow dedup registry: set of (from_name, to_name) tuples
        # Prevents duplicate arrows between the same pair of icons
        self._connections = set()

        # Step 1: Create Presentation & Blank Slide
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.50)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Step 2: Add Title
        self._add_title(title)

    def _add_title(self, title):
        """Step 2: Add diagram title."""
        title_box = self.slide.shapes.add_textbox(
            Inches(0.26), Inches(0.40), Inches(12.81), Inches(0.70)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(20)
        run.font.bold = True

    # ─── Step 3: Groups ─────────────────────────────────────────────────────

    # Group icon occupies 0.42×0.42" at top-left corner.
    # Service icons inside a group should be offset to avoid overlap:
    #   safe_x >= group_left + 0.55"  (clear of group icon)
    #   safe_y >= group_top + 0.55"   (clear of group label)
    GROUP_ICON_CLEARANCE = 0.55  # inches to clear group icon + label

    @staticmethod
    def safe_icon_position(group_left, group_top, group_width, group_height,
                           col=0, row=0, cols=1, icon_type='service'):
        """
        Calculate a safe (center_x, top_y) for placing an icon inside a group,
        avoiding overlap with the group corner icon (0.42") and label.

        Args:
            group_left, group_top: Group position in inches
            group_width, group_height: Group size in inches
            col, row: Grid position (0-indexed)
            cols: Total columns in row
            icon_type: 'service' (0.83") or 'resource' (0.50")

        Returns:
            (center_x_in, top_y_in) tuple
        """
        clearance = AWSArchitectureDiagram.GROUP_ICON_CLEARANCE
        icon_w = 0.83 if icon_type == 'service' else 0.50
        usable_left = group_left + clearance
        usable_width = group_width - clearance - 0.2  # right margin
        usable_top = group_top + clearance
        usable_height = group_height - clearance - 0.2

        cell_w = usable_width / max(cols, 1)
        center_x = usable_left + cell_w * col + cell_w / 2
        top_y = usable_top + row * (icon_w + 0.6)  # 0.6" for label space

        return center_x, top_y

    def add_group(self, group_type, left_in, top_in, width_in, height_in,
                  label=None, custom_color=None, custom_dash=None):
        """
        Step 3: Add a group container.

        Args:
            group_type: Key from GROUP_SPECS (e.g., "AWS Cloud", "VPC")
            left_in, top_in: Position in inches
            width_in, height_in: Size in inches
            label: Override label text (default: group_type name)
            custom_color: Override border RGBColor
            custom_dash: Override dash style

        Returns:
            The rectangle shape object
        """
        spec = GROUP_SPECS.get(group_type, GROUP_SPECS["Generic"])
        color = custom_color or spec["color"]
        dash = custom_dash if custom_dash is not None else spec["dash"]
        label_text = label or group_type

        left = Inches(left_in)
        top = Inches(top_in)
        width = Inches(width_in)
        height = Inches(height_in)

        # Rectangle container
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.background()  # Transparent
        shape.line.width = Pt(1.25)
        shape.line.color.rgb = color
        if dash:
            shape.line.dash_style = dash

        # Group label as separate text box (top-left, after icon space)
        # Position: offset inside border to avoid overlap with group icon and border
        # Label width: just enough for text, NOT spanning entire group
        label_left = left + Inches(0.48)  # after group icon (0.42")
        label_top = top + Inches(0.06)    # inside border, avoids overlap with border line
        max_label_w = min(width_in - 0.6, 2.50)  # cap at 2.50" to avoid overlap with icons
        label_box = self.slide.shapes.add_textbox(
            label_left, label_top, Inches(max_label_w), Inches(0.22)
        )
        tf = label_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = label_text
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black per AWS rules

        # Clear the rectangle's own text frame (was causing faint text on border)
        shape.text_frame.paragraphs[0].text = ""

        # Group icon (top-left corner, 0.42in)
        icon_file = spec.get("icon")
        if icon_file:
            icon_path = os.path.join(self.group_icon_dir, icon_file)
            if os.path.exists(icon_path):
                self.slide.shapes.add_picture(
                    icon_path, left, top, Inches(0.42), Inches(0.42)
                )

        return shape

    # ─── Step 4: Service Icons ───────────────────────────────────────────────

    def add_service_icon(self, icon_path, label_text, center_x_in, top_y_in, name=None):
        """
        Step 4: Add a service icon with label below.

        Args:
            icon_path: Path to the 48px PNG icon file
            label_text: Label text (e.g., "AWS Lambda")
            center_x_in: Center X position in inches
            top_y_in: Top Y position in inches
            name: Optional unique name for arrow anchoring via connect().
                  If not provided, uses label_text (first line).

        Returns:
            Tuple of (picture_shape, textbox_shape)
        """
        s = self.scale
        icon_w = Inches(self.SERVICE_ICON_SIZE * s)
        icon_h = Inches(self.SERVICE_ICON_SIZE * s)
        label_w_in = self.SERVICE_LABEL_W * s
        label_w = Inches(label_w_in)
        label_h = Inches(0.40)
        font_size = 12  # AWS official: 12pt Arial Regular
        center_x = Inches(center_x_in)
        top_y = Inches(top_y_in)

        # Icon image
        pic_left = center_x - icon_w // 2
        pic = self.slide.shapes.add_picture(icon_path, pic_left, top_y, icon_w, icon_h)

        # Label below — 12pt Arial Regular per AWS rules
        label_left = center_x - label_w // 2
        label_top = top_y + icon_h
        txBox = self.slide.shapes.add_textbox(label_left, label_top, label_w, label_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label_text
        p.alignment = PP_ALIGN.CENTER
        # Apply font to ALL runs (word wrap may create multiple runs)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Register icon for connect()
        icon_name = name or label_text.split('\n')[0].strip()
        icon_size_in = self.SERVICE_ICON_SIZE * s
        label_h_in = 0.40
        self._icons[icon_name] = {
            'cx': center_x_in,
            'cy': top_y_in + icon_size_in / 2,
            'w': icon_size_in,
            'h': icon_size_in,
            'type': 'service',
            'label_bottom': top_y_in + icon_size_in + label_h_in,
            'label_w': label_w_in,
        }

        return pic, txBox

    # ─── Step 5: Resource Icons ──────────────────────────────────────────────

    def add_resource_icon(self, icon_path, label_text, center_x_in, top_y_in, name=None):
        """
        Step 5: Add a resource icon (0.50×0.50in) with label below.

        Args:
            icon_path: Path to the resource PNG icon file
            label_text: Label text (e.g., "Instance")
            center_x_in: Center X position in inches
            top_y_in: Top Y position in inches
            name: Optional unique name for arrow anchoring via connect()

        Returns:
            Tuple of (picture_shape, textbox_shape)
        """
        s = self.scale
        icon_w = Inches(self.RESOURCE_ICON_SIZE * s)
        icon_h = Inches(self.RESOURCE_ICON_SIZE * s)
        label_w_in = self.RESOURCE_LABEL_W * s
        label_w = Inches(label_w_in)
        label_h = Inches(0.30)
        font_size = 12  # AWS official: 12pt Arial Regular
        center_x = Inches(center_x_in)
        top_y = Inches(top_y_in)

        pic_left = center_x - icon_w // 2
        pic = self.slide.shapes.add_picture(icon_path, pic_left, top_y, icon_w, icon_h)

        label_left = center_x - label_w // 2
        label_top = top_y + icon_h
        txBox = self.slide.shapes.add_textbox(label_left, label_top, label_w, label_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label_text
        p.alignment = PP_ALIGN.CENTER
        # Apply font to ALL runs (word wrap may create multiple runs)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Register icon for connect()
        icon_name = name or label_text.split('\n')[0].strip()
        icon_size_in = self.RESOURCE_ICON_SIZE * s
        label_h_in = 0.30
        self._icons[icon_name] = {
            'cx': center_x_in,
            'cy': top_y_in + icon_size_in / 2,
            'w': icon_size_in,
            'h': icon_size_in,
            'type': 'resource',
            'label_bottom': top_y_in + icon_size_in + label_h_in,
            'label_w': label_w_in,
        }

        return pic, txBox

    # ─── Step 6: Arrows & Connectors ─────────────────────────────────────────

    def _icon_edge_point(self, icon_name, direction):
        """
        Get the edge midpoint of a registered icon.

        Args:
            icon_name: Name used when registering icon
            direction: 'top', 'bottom', 'left', 'right'

        Returns:
            (x_in, y_in) at the edge midpoint
        """
        info = self._icons[icon_name]
        cx, cy = info['cx'], info['cy']
        hw, hh = info['w'] / 2, info['h'] / 2

        if direction == 'right':
            return cx + hw, cy
        elif direction == 'left':
            return cx - hw, cy
        elif direction == 'top':
            return cx, cy - hh
        elif direction == 'bottom':
            # Start below label to avoid arrow overlapping label text
            label_bottom = info.get('label_bottom', cy + hh)
            return cx, label_bottom
        return cx, cy

    def _auto_direction(self, from_name, to_name):
        """Determine best exit/entry directions between two icons."""
        f = self._icons[from_name]
        t = self._icons[to_name]
        dx = t['cx'] - f['cx']
        dy = t['cy'] - f['cy']

        if abs(dx) >= abs(dy):
            return ('right', 'left') if dx > 0 else ('left', 'right')
        else:
            return ('bottom', 'top') if dy > 0 else ('top', 'bottom')

    def _segment_crosses_icon(self, x1, y1, x2, y2, exclude_names=None, check_labels=True):
        """
        Check if a horizontal or vertical line segment crosses through any
        registered icon OR its label area (excluding the named icons).

        Args:
            x1, y1, x2, y2: Segment endpoints
            exclude_names: Icon names to skip (source/target)
            check_labels: Also check label text areas below icons

        Returns list of icon names that the segment crosses through.
        """
        exclude = set(exclude_names or [])
        crossed = []
        margin = 0.05  # small margin to avoid false positives at edges

        for name, info in self._icons.items():
            if name in exclude:
                continue
            cx, cy = info['cx'], info['cy']
            hw, hh = info['w'] / 2 + margin, info['h'] / 2 + margin
            icon_left = cx - hw
            icon_right = cx + hw
            icon_top = cy - hh
            icon_bottom = cy + hh

            # Extend bottom to include label area if check_labels=True
            if check_labels and 'label_bottom' in info:
                label_bottom = info['label_bottom'] + margin
                label_hw = info.get('label_w', info['w']) / 2 + margin
                label_left = cx - label_hw
                label_right = cx + label_hw
            else:
                label_bottom = icon_bottom
                label_left = icon_left
                label_right = icon_right

            # Combined bounding box (icon + label)
            full_left = min(icon_left, label_left)
            full_right = max(icon_right, label_right)
            full_top = icon_top
            full_bottom = label_bottom

            # Horizontal segment (y1 ≈ y2)
            if abs(y1 - y2) < 0.01:
                y = y1
                seg_left, seg_right = min(x1, x2), max(x1, x2)
                if full_top < y < full_bottom and seg_left < full_right and seg_right > full_left:
                    crossed.append(name)
            # Vertical segment (x1 ≈ x2)
            elif abs(x1 - x2) < 0.01:
                x = x1
                seg_top, seg_bottom = min(y1, y2), max(y1, y2)
                if full_left < x < full_right and seg_top < full_bottom and seg_bottom > full_top:
                    crossed.append(name)

        return crossed

    def connect_routed(self, from_name, to_name, waypoints,
                       from_dir=None, to_dir=None, bidirectional=False):
        """
        Connect two icons through explicit waypoints (for routing around obstacles).
        All segments are horizontal or vertical. Waypoints are (x, y) tuples.

        Args:
            from_name: Source icon name
            to_name: Target icon name
            waypoints: List of (x, y) intermediate points
            from_dir: Exit direction from source
            to_dir: Entry direction to target
            bidirectional: Arrow heads on both ends

        Example:
            # Route EC2→DynamoDB going below private subnet
            d.connect_routed('EC2', 'DDB', [(7.0, 6.5), (12.0, 6.5)],
                             from_dir='bottom', to_dir='bottom')
        """
        if from_name not in self._icons or to_name not in self._icons:
            raise ValueError("Icon not registered")

        if from_dir is None or to_dir is None:
            auto_from, auto_to = self._auto_direction(from_name, to_name)
            from_dir = from_dir or auto_from
            to_dir = to_dir or auto_to

        sx, sy = self._icon_edge_point(from_name, from_dir)
        ex, ey = self._icon_edge_point(to_name, to_dir)

        # Build full path: start → waypoints → end
        points = [(sx, sy)] + list(waypoints) + [(ex, ey)]

        # Draw elbow connectors between consecutive points
        shapes = []
        for i in range(len(points) - 1):
            p1 = points[i]
            p2 = points[i + 1]
            is_last = (i == len(points) - 2)
            is_first = (i == 0)

            # If axis-aligned, use straight arrow/line
            if abs(p1[0] - p2[0]) < 0.05 or abs(p1[1] - p2[1]) < 0.05:
                if is_last:
                    seg = self.add_arrow(p1[0], p1[1], p2[0], p2[1],
                                         bidirectional=False)
                else:
                    seg = self.add_line(p1[0], p1[1], p2[0], p2[1])
                if bidirectional and is_first:
                    self._add_arrowhead(seg, end='head')
            else:
                # L-shaped: use elbow connector
                bidir = bidirectional if is_first else False
                has_arrow = is_last
                if has_arrow:
                    seg = self.add_elbow_arrow(p1[0], p1[1], p2[0], p2[1],
                                               bidirectional=bidir)
                else:
                    seg = self.add_elbow_arrow(p1[0], p1[1], p2[0], p2[1],
                                               bidirectional=bidir)
            shapes.append(seg)

        # Register connection for dedup
        pair = (min(from_name, to_name), max(from_name, to_name))
        self._connections.add(pair)

        return shapes

    def connect(self, from_name, to_name, from_dir=None, to_dir=None,
                bidirectional=False):
        """
        Connect two registered icons with an orthogonal (right-angle) arrow.
        ALL arrows are solid and routed orthogonally — no diagonal lines,
        no dashed lines (per AWS Architecture Icons PPTX deck slide 16).
        If icons are axis-aligned, a straight horizontal/vertical arrow is used.
        Duplicate connections between the same pair are silently skipped.

        Args:
            from_name: Source icon name (as registered)
            to_name: Target icon name (as registered)
            from_dir: Exit direction ('top','bottom','left','right') or None=auto
            to_dir: Entry direction or None=auto
            bidirectional: Arrow heads on both ends

        Returns:
            Connector shape(s)
        """
        if from_name not in self._icons:
            raise ValueError(f"Icon '{from_name}' not registered. Available: {list(self._icons.keys())}")
        if to_name not in self._icons:
            raise ValueError(f"Icon '{to_name}' not registered. Available: {list(self._icons.keys())}")

        # Dedup: skip if this connection already exists
        pair = (min(from_name, to_name), max(from_name, to_name))
        if pair in self._connections:
            return None
        self._connections.add(pair)

        # Always use orthogonal routing
        return self.connect_orthogonal(from_name, to_name, from_dir, to_dir,
                                       bidirectional=bidirectional)

    def connect_orthogonal(self, from_name, to_name, from_dir=None, to_dir=None,
                           bidirectional=False):
        """
        Connect two icons with an elbow (right-angle) connector.
        Uses PowerPoint's native elbow connector for clean right-angle routing.
        All arrows are SOLID per AWS Architecture Icons guidelines.

        Args:
            from_name, to_name: Registered icon names
            from_dir, to_dir: Exit/entry directions or None=auto
            bidirectional: Arrow heads on both ends
        """
        if from_name not in self._icons or to_name not in self._icons:
            raise ValueError(f"Icon not registered")

        if from_dir is None or to_dir is None:
            auto_from, auto_to = self._auto_direction(from_name, to_name)
            from_dir = from_dir or auto_from
            to_dir = to_dir or auto_to

        sx, sy = self._icon_edge_point(from_name, from_dir)
        ex, ey = self._icon_edge_point(to_name, to_dir)

        # If icons are aligned on the axis, use a straight arrow
        if abs(sy - ey) < 0.05 or abs(sx - ex) < 0.05:
            return self.add_arrow(sx, sy, ex, ey, bidirectional=bidirectional), None

        # Use PowerPoint elbow connector for L-shaped routing
        connector = self.add_elbow_arrow(sx, sy, ex, ey, bidirectional=bidirectional)
        return connector, None

    def add_arrow(self, start_x_in, start_y_in, end_x_in, end_y_in,
                  bidirectional=False):
        """
        Step 6: Add a straight arrow connector. ALL arrows are SOLID.

        Args:
            start_x_in, start_y_in: Start position in inches
            end_x_in, end_y_in: End position in inches
            bidirectional: If True, add arrow heads on both ends

        Returns:
            The connector shape object
        """
        connector = self.slide.shapes.add_connector(
            1,  # MSO_CONNECTOR_TYPE.STRAIGHT
            Inches(start_x_in), Inches(start_y_in),
            Inches(end_x_in), Inches(end_y_in),
        )
        connector.line.width = Pt(1.25)

        # Add arrow head(s) via XML manipulation
        sp_pr = connector._element.find(qn('p:cxnSp'))
        if sp_pr is None:
            sp_pr = connector._element
        ln = sp_pr.find('.//' + qn('a:ln'))
        if ln is None:
            # Create line element
            sp_pr_elem = sp_pr.find(qn('p:spPr'))
            if sp_pr_elem is None:
                sp_pr_elem = sp_pr.find(qn('xdr:spPr'))
            if sp_pr_elem is None:
                sp_pr_elem = sp_pr
            from lxml import etree
            ln = etree.SubElement(sp_pr_elem, qn('a:ln'))

        # Tail arrow (end)
        tail = ln.makeelement(qn('a:tailEnd'), {})
        tail.set('type', 'arrow')
        tail.set('w', 'lg')
        tail.set('len', 'lg')
        ln.append(tail)

        # Head arrow (start) for bidirectional
        if bidirectional:
            head = ln.makeelement(qn('a:headEnd'), {})
            head.set('type', 'arrow')
            head.set('w', 'lg')
            head.set('len', 'lg')
            ln.append(head)

        return connector

    def add_elbow_arrow(self, start_x_in, start_y_in, end_x_in, end_y_in,
                        bidirectional=False):
        """
        Add an elbow (right-angle) connector with arrowhead.
        Uses PowerPoint's native elbow connector type for clean L-shaped routing.
        Users can adjust the bend point by dragging the yellow handle in PPT.

        Args:
            start_x_in, start_y_in: Start position in inches
            end_x_in, end_y_in: End position in inches
            bidirectional: If True, add arrow heads on both ends

        Returns:
            The connector shape object
        """
        connector = self.slide.shapes.add_connector(
            2,  # MSO_CONNECTOR_TYPE.ELBOW (native right-angle routing)
            Inches(start_x_in), Inches(start_y_in),
            Inches(end_x_in), Inches(end_y_in),
        )
        connector.line.width = Pt(1.25)

        # Add arrow head(s)
        sp_pr = connector._element.find(qn('p:cxnSp'))
        if sp_pr is None:
            sp_pr = connector._element
        ln = sp_pr.find('.//' + qn('a:ln'))
        if ln is None:
            sp_pr_elem = sp_pr.find(qn('p:spPr'))
            if sp_pr_elem is None:
                sp_pr_elem = sp_pr.find(qn('xdr:spPr'))
            if sp_pr_elem is None:
                sp_pr_elem = sp_pr
            from lxml import etree
            ln = etree.SubElement(sp_pr_elem, qn('a:ln'))

        tail = ln.makeelement(qn('a:tailEnd'), {})
        tail.set('type', 'arrow')
        tail.set('w', 'lg')
        tail.set('len', 'lg')
        ln.append(tail)

        if bidirectional:
            head = ln.makeelement(qn('a:headEnd'), {})
            head.set('type', 'arrow')
            head.set('w', 'lg')
            head.set('len', 'lg')
            ln.append(head)

        return connector

    def add_line(self, start_x_in, start_y_in, end_x_in, end_y_in):
        """Add a plain line (no arrow heads). ALL lines are SOLID."""
        connector = self.slide.shapes.add_connector(
            1, Inches(start_x_in), Inches(start_y_in),
            Inches(end_x_in), Inches(end_y_in),
        )
        connector.line.width = Pt(1.25)
        return connector

    def _add_arrowhead(self, connector, end='tail'):
        """Add an arrowhead to an existing connector."""
        sp_pr = connector._element.find(qn('p:cxnSp'))
        if sp_pr is None:
            sp_pr = connector._element
        ln = sp_pr.find('.//' + qn('a:ln'))
        if ln is not None:
            tag = qn(f'a:{end}End')
            elem = ln.makeelement(tag, {})
            elem.set('type', 'arrow')
            elem.set('w', 'lg')
            elem.set('len', 'lg')
            ln.append(elem)

    # ─── Step 7: Annotations ────────────────────────────────────────────────

    def add_text_annotation(self, text, left_in, top_in, width_in=2.0, height_in=0.3,
                            font_size=12, bold=False, alignment=PP_ALIGN.LEFT):
        """Step 7: Add a text annotation."""
        txBox = self.slide.shapes.add_textbox(
            Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        return txBox

    def add_numbered_step(self, number, center_x_in, center_y_in):
        """Add a numbered step oval (0.36×0.36in)."""
        from pptx.enum.shapes import MSO_SHAPE
        size = Inches(0.36)
        left = Inches(center_x_in) - size // 2
        top = Inches(center_y_in) - size // 2
        shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.background()
        shape.line.width = Pt(1.0)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(10)
        run.font.bold = True
        return shape

    def add_footer(self):
        """Add standard AWS copyright footer."""
        self.add_text_annotation(
            "© 2026, Amazon Web Services, Inc. or its affiliates.",
            4.42, 6.88, 4.50, 0.40, font_size=8
        )

    def add_general_icon(self, icon_path, label_text, center_x_in, top_y_in, name=None):
        """
        Add a general icon (0.50×0.50in) with label below.
        General icons represent non-service elements: users, internet, servers, etc.
        See assets/general-icons/ for available icons.

        Args:
            icon_path: Path to the general icon PNG file
            label_text: Label text (e.g., "Users", "Internet")
            center_x_in: Center X position in inches
            top_y_in: Top Y position in inches
            name: Optional unique name for arrow anchoring via connect()

        Returns:
            Tuple of (picture_shape, textbox_shape)
        """
        return self.add_resource_icon(icon_path, label_text, center_x_in, top_y_in, name=name)

    def add_legend(self, left_in=10.0, top_in=6.5, width_in=3.0, height_in=0.7):
        """Add a legend box explaining arrow conventions (all solid arrows per AWS guidelines)."""
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_in), Inches(top_in),
            Inches(width_in), Inches(height_in)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.width = Pt(0.5)
        shape.line.color.rgb = RGBColor(0x54, 0x5B, 0x64)
        shape.text_frame.paragraphs[0].text = ""

        txBox = self.slide.shapes.add_textbox(
            Inches(left_in + 0.15), Inches(top_in + 0.1),
            Inches(width_in - 0.3), Inches(height_in - 0.2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "→  Data / request flow"
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x23, 0x2F, 0x3E)
        p2 = tf.add_paragraph()
        p2.text = "⇄  Bidirectional flow"
        p2.space_before = Pt(4)
        run2 = p2.runs[0]
        run2.font.name = "Arial"
        run2.font.size = Pt(8)
        run2.font.color.rgb = RGBColor(0x23, 0x2F, 0x3E)

    # ─── Step 8: Save ────────────────────────────────────────────────────────

    def save(self, output_path):
        """Step 8: Save the presentation."""
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        self.prs.save(output_path)
        print(f"Diagram saved to: {output_path}")
        return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# Example Usage
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    # Example: Basic VPC Architecture
    diagram = AWSArchitectureDiagram("Example: Basic VPC Architecture")

    # Step 3: Groups (outermost → innermost)
    diagram.add_group("AWS Cloud", 1.0, 1.2, 11.5, 5.5)
    diagram.add_group("VPC", 1.5, 1.8, 10.5, 4.5, label="VPC (10.0.0.0/16)")
    diagram.add_group("Availability Zone", 2.0, 2.4, 4.5, 3.5, label="Availability Zone 1")
    diagram.add_group("Public Subnet", 2.2, 2.8, 4.1, 1.3, label="Public Subnet")
    diagram.add_group("Private Subnet", 2.2, 4.3, 4.1, 1.3, label="Private Subnet")
    diagram.add_group("Availability Zone", 7.0, 2.4, 4.5, 3.5, label="Availability Zone 2")
    diagram.add_group("Public Subnet", 7.2, 2.8, 4.1, 1.3, label="Public Subnet")
    diagram.add_group("Private Subnet", 7.2, 4.3, 4.1, 1.3, label="Private Subnet")

    # Step 7: Annotations
    diagram.add_footer()

    # Step 8: Save
    diagram.save("example_vpc_architecture.pptx")
