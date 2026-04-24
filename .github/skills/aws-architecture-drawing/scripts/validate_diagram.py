#!/usr/bin/env python3
"""
Validate an AWS architecture PPTX diagram against design rules.

Checks performed:
  1. Slide dimensions (13.33 × 7.50 in widescreen)
  2. Non-elbow connectors detected (all L-shaped connections must use elbow type)
  3. Text size vs icon size (label width and font height must not exceed icon)
  4. Overlap detection: icon-icon, label-label, icon-label, arrow-icon, arrow-label
  5. Nested group containment (child groups must be fully inside parent groups)
  6. Boundary checks (shapes must stay within slide bounds)
  7. Font size range validation
  8. AWS group color validation
  9. Arrow style validation (must be solid, not dashed)
  10. Group proportionality (aspect ratio, minimum size, balanced dimensions)
  11. Group padding (icons must have minimum distance from group borders)
  12. Icon spacing (minimum distance between adjacent icons)
"""

import argparse
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ─── Constants ────────────────────────────────────────────────────────────────

SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.50
SLIDE_W = Inches(SLIDE_W_IN)
SLIDE_H = Inches(SLIDE_H_IN)
EMU_PER_INCH = 914400
MIN_FONT_PT = 6
MAX_FONT_PT = 28
OVERLAP_TOLERANCE_IN = 0.02  # allow 0.02" overlap tolerance

GROUP_COLORS = {
    '232f3e': 'AWS Cloud',
    '8c4fff': 'VPC / Generic Purple',
    '00a4a6': 'Region / AZ / Private Subnet',
    '7aa116': 'Public Subnet',
    'dd344c': 'Security Group',
    'ed7100': 'Auto Scaling',
    'e7157b': 'AWS Account',
    '01a88d': 'Corporate DC / Server Contents',
    'c925d1': 'Spot Fleet',
    '7d8998': 'Generic / External',
}

APPROVED_BORDER_COLORS = set(GROUP_COLORS.keys())


# ─── Helpers ──────────────────────────────────────────────────────────────────

def emu_to_in(emu):
    """Convert EMU to inches."""
    return emu / EMU_PER_INCH


def rect_of(shape):
    """Return (left, top, right, bottom) in inches."""
    l = emu_to_in(shape.left)
    t = emu_to_in(shape.top)
    r = l + emu_to_in(shape.width)
    b = t + emu_to_in(shape.height)
    return (l, t, r, b)


def rects_overlap(a, b, tol=OVERLAP_TOLERANCE_IN):
    """Check if two rectangles overlap (beyond tolerance)."""
    al, at, ar, ab = a
    bl, bt, br, bb = b
    # No overlap if separated
    if ar - tol <= bl or br - tol <= al:
        return False
    if ab - tol <= bt or bb - tol <= at:
        return False
    return True


def rect_contains(outer, inner, tol=OVERLAP_TOLERANCE_IN):
    """Check if outer rectangle fully contains inner rectangle."""
    ol, ot, oright, ob = outer
    il, it, iright, ib = inner
    return (il >= ol - tol and it >= ot - tol and
            iright <= oright + tol and ib <= ob + tol)


def rect_str(r):
    return f'({r[0]:.2f}", {r[1]:.2f}") to ({r[2]:.2f}", {r[3]:.2f}")'


# ─── Validator ────────────────────────────────────────────────────────────────

class DiagramValidator:
    def __init__(self, filepath):
        self.filepath = filepath
        self.prs = Presentation(filepath)
        self.errors = []
        self.warnings = []
        self.info = []

        # Collected shape data for cross-checks
        self.icons = []       # list of (name, rect, shape)
        self.labels = []      # list of (text, rect, shape)
        self.groups = []      # list of (name, rect, shape)
        self.connectors = []  # list of (type_name, shape)

    def validate(self):
        """Run all validation checks."""
        self._check_slide_dimensions()

        for slide in self.prs.slides:
            self._collect_shapes(slide)

        self._check_non_elbow_connectors()
        self._check_label_vs_icon_size()
        self._check_overlaps()
        self._check_nested_group_containment()
        self._check_boundary()
        self._check_arrow_styles()
        self._check_group_proportionality()
        self._check_group_padding()
        self._check_icon_spacing()

        return len(self.errors) == 0

    # ─── Step 1: Slide Dimensions ─────────────────────────────────────────

    def _check_slide_dimensions(self):
        w = emu_to_in(self.prs.slide_width)
        h = emu_to_in(self.prs.slide_height)
        self.info.append(f'Slide dimensions: {w:.2f}" x {h:.2f}"')
        self.info.append(f'Total slides: {len(self.prs.slides)}')

        if abs(w - SLIDE_W_IN) > 0.1:
            self.errors.append(f'Slide width {w:.2f}" != expected {SLIDE_W_IN}"')
        if abs(h - SLIDE_H_IN) > 0.1:
            self.errors.append(f'Slide height {h:.2f}" != expected {SLIDE_H_IN}"')
        if len(self.prs.slides) == 0:
            self.errors.append('Presentation has no slides')

    # ─── Collect shapes ───────────────────────────────────────────────────

    def _collect_shapes(self, slide):
        icon_count = 0
        group_count = 0
        connector_count = 0
        text_count = 0

        for shape in slide.shapes:
            # Icons (pictures) — skip group corner icons (< 0.35")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                r = rect_of(shape)
                icon_w = r[2] - r[0]
                icon_h = r[3] - r[1]
                if icon_w < 0.45 and icon_h < 0.45:
                    continue  # group corner icon, skip
                icon_count += 1
                name = f'icon@({r[0]:.1f},{r[1]:.1f})'
                self.icons.append((name, r, shape))

            # Connectors
            elif hasattr(shape, 'begin_x'):
                connector_count += 1
                ctype = self._get_connector_type(shape)
                self.connectors.append((ctype, shape))

            # Auto shapes (groups = rectangles with borders)
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if hasattr(shape, 'line') and shape.line.fill.type is not None:
                    group_count += 1
                    r = rect_of(shape)
                    color_hex = ''
                    if shape.line.color and shape.line.color.rgb:
                        color_hex = str(shape.line.color.rgb).lower()
                    gname = GROUP_COLORS.get(color_hex, f'Group#{color_hex}')
                    self.groups.append((gname, r, shape))
                    self._check_group_color(color_hex, r)

            # Text boxes — skip empty or very short text
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if len(text) < 2:
                    continue  # skip empty/trivial text
                text_count += 1
                r = rect_of(shape)
                display = text[:40].replace('\n', ' ')
                self.labels.append((display, r, shape))
                self._check_font(shape)

        self.info.append(
            f'Shapes: {icon_count} icons, {group_count} groups, '
            f'{connector_count} connectors, {text_count} text elements'
        )

    def _get_connector_type(self, shape):
        """Get connector type: 'straight', 'elbow', or 'curve'."""
        elem = shape._element
        # Look for cxnSp/spPr/prstGeom
        prst = elem.find('.//' + qn('a:prstGeom'))
        if prst is not None:
            prst_name = prst.get('prst', '')
            if 'bentConnector' in prst_name:
                return 'elbow'
            elif 'straightConnector' in prst_name:
                return 'straight'
            elif 'curvedConnector' in prst_name:
                return 'curve'
        return 'unknown'

    # ─── Check 2: Non-elbow connectors ────────────────────────────────────

    def _check_non_elbow_connectors(self):
        straight_count = 0
        elbow_count = 0
        for ctype, shape in self.connectors:
            if ctype == 'straight':
                # Check if it's axis-aligned (horizontal or vertical) — those are OK
                try:
                    bx = emu_to_in(shape.begin_x)
                    by = emu_to_in(shape.begin_y)
                    ex = emu_to_in(shape.end_x)
                    ey = emu_to_in(shape.end_y)
                    if abs(bx - ex) < 0.05 or abs(by - ey) < 0.05:
                        # Axis-aligned straight connector is fine
                        continue
                    else:
                        # Diagonal straight connector — should be elbow
                        straight_count += 1
                        self.errors.append(
                            f'Non-elbow diagonal connector found from '
                            f'({bx:.2f}", {by:.2f}") to ({ex:.2f}", {ey:.2f}"). '
                            f'Use elbow connector (type=2) for L-shaped routing'
                        )
                except Exception:
                    straight_count += 1
            elif ctype == 'elbow':
                elbow_count += 1

        if elbow_count > 0:
            self.info.append(f'Elbow connectors: {elbow_count}')
        if straight_count == 0 and elbow_count > 0:
            self.info.append('All L-shaped connections use elbow connectors ✓')

    # ─── Check 3: Text size vs icon size ──────────────────────────────────

    def _check_label_vs_icon_size(self):
        """Check that labels are not wider than their associated icons."""
        for ltext, lrect, lshape in self.labels:
            # Find the closest icon above this label
            ll, lt, lr, lb = lrect
            label_w = lr - ll
            label_cx = (ll + lr) / 2

            best_icon = None
            best_dist = float('inf')
            for iname, irect, ishape in self.icons:
                il, it, ir, ib = irect
                icon_cx = (il + ir) / 2
                icon_w = ir - il
                # Label should be directly below icon (within 0.3" horizontal)
                if abs(label_cx - icon_cx) < 0.3 and lt >= ib - 0.05:
                    dist = lt - ib
                    if dist < best_dist and dist < 0.5:
                        best_dist = dist
                        best_icon = (iname, irect, ishape, icon_w)

            if best_icon:
                iname, irect, ishape, icon_w = best_icon
                if label_w > icon_w + 0.40:  # allow label up to 0.40" wider than icon
                    self.errors.append(
                        f'Label "{ltext}" width {label_w:.2f}" exceeds icon '
                        f'width {icon_w:.2f}" at {iname}. '
                        f'Label must not be wider than its icon'
                    )

                # Check font point size vs icon size in points
                # Icon of 0.7" ≈ 50pt height. Font should be well below that.
                icon_h_pt = icon_w * 72  # convert inches to points
                for para in lshape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_pt = run.font.size / 12700
                            if font_pt > icon_h_pt * 0.3:  # font > 30% of icon size
                                self.errors.append(
                                    f'Label "{ltext}" font {font_pt:.0f}pt is too large '
                                    f'relative to icon ({icon_w:.2f}" = {icon_h_pt:.0f}pt). '
                                    f'Font should be ≤{icon_h_pt*0.3:.0f}pt'
                                )

    # ─── Check 4: Overlap detection ───────────────────────────────────────

    def _check_overlaps(self):
        # Icon-icon overlap
        for i in range(len(self.icons)):
            for j in range(i + 1, len(self.icons)):
                n1, r1, _ = self.icons[i]
                n2, r2, _ = self.icons[j]
                if rects_overlap(r1, r2):
                    self.errors.append(
                        f'Icon overlap: {n1} {rect_str(r1)} overlaps {n2} {rect_str(r2)}'
                    )

        # Categorize labels
        title_labels = set()  # indices of title/group-header labels
        icon_labels = set()   # indices of icon labels
        group_labels = set()  # indices of group/subnet labels

        for idx, (ltext, lrect, _) in enumerate(self.labels):
            lw = lrect[2] - lrect[0]
            # Title labels are very wide (>8")
            if lw > 8.0:
                title_labels.add(idx)
                continue
            # Check if this label belongs to a group (overlaps group border area)
            is_group_label = False
            for gname, grect, _ in self.groups:
                gl, gt, gr, gb = grect
                lcx = (lrect[0] + lrect[2]) / 2
                lcy = (lrect[1] + lrect[3]) / 2
                # Group header label: near top of group
                if gl <= lcx <= gr and abs(lrect[1] - gt) < 0.5:
                    is_group_label = True
                    break
            if is_group_label:
                group_labels.add(idx)
            else:
                icon_labels.add(idx)

        # Label-label overlap: icon labels vs each other
        for i in icon_labels:
            for j in icon_labels:
                if j <= i:
                    continue
                t1, r1, _ = self.labels[i]
                t2, r2, _ = self.labels[j]
                if rects_overlap(r1, r2, tol=0.05):
                    self.warnings.append(
                        f'Label overlap: "{t1}" and "{t2}" overlap at {rect_str(r1)}'
                    )

        # Group-label vs icon-label overlap (e.g., wide VPC label overlapping DynamoDB label)
        for gi in group_labels:
            gt_text, gr_rect, _ = self.labels[gi]
            gr_w = gr_rect[2] - gr_rect[0]
            for ii in icon_labels:
                it_text, ir_rect, _ = self.labels[ii]
                if rects_overlap(gr_rect, ir_rect, tol=0.05):
                    self.warnings.append(
                        f'Group label "{gt_text}" ({gr_w:.1f}" wide) overlaps '
                        f'icon label "{it_text}"'
                    )

        # Group-label vs icon overlap (group label textbox overlapping a service icon)
        for gi in group_labels:
            gt_text, gr_rect, _ = self.labels[gi]
            for iname, irect, _ in self.icons:
                if rects_overlap(gr_rect, irect, tol=0.05):
                    self.warnings.append(
                        f'Group label "{gt_text}" overlaps {iname}'
                    )

        # Icon-label overlap: only check against icon labels (not group/title labels)
        for iname, irect, _ in self.icons:
            for idx in icon_labels:
                ltext, lrect, _ = self.labels[idx]
                if rects_overlap(irect, lrect, tol=0.03):
                    # Check if this label belongs to this icon (directly below)
                    icx = (irect[0] + irect[2]) / 2
                    lcx = (lrect[0] + lrect[2]) / 2
                    if abs(icx - lcx) < 0.3 and lrect[1] >= irect[3] - 0.05:
                        continue  # label belongs to this icon
                    self.warnings.append(
                        f'Icon {iname} overlaps label "{ltext}"'
                    )

        # Arrow-icon overlap: check if any connector passes through an icon
        self._check_arrow_overlaps()

    def _check_arrow_overlaps(self):
        """Check if arrows pass through icons or labels (not their endpoints)."""
        for ctype, shape in self.connectors:
            try:
                bx = emu_to_in(shape.begin_x)
                by = emu_to_in(shape.begin_y)
                ex = emu_to_in(shape.end_x)
                ey = emu_to_in(shape.end_y)
            except Exception:
                continue

            # Build arrow bounding box with some padding
            al = min(bx, ex)
            at = min(by, ey)
            ar = max(bx, ex)
            ab = max(by, ey)
            arrow_rect = (al, at, ar, ab)

            # Check against icons
            for iname, irect, _ in self.icons:
                il, it, ir, ib = irect
                icx = (il + ir) / 2
                icy = (it + ib) / 2

                # Skip if arrow starts or ends at this icon (expected)
                if (abs(bx - icx) < (ir - il) and abs(by - icy) < (ib - it)):
                    continue
                if (abs(ex - icx) < (ir - il) and abs(ey - icy) < (ib - it)):
                    continue

                # Check if arrow line segment intersects icon rectangle
                if self._line_intersects_rect(bx, by, ex, ey, irect, tol=0.05):
                    self.errors.append(
                        f'Arrow ({bx:.2f},{by:.2f})→({ex:.2f},{ey:.2f}) '
                        f'passes through {iname}. Arrows must not cross icons'
                    )

            # Check against icon labels only (arrows naturally pass through group/title areas)
            for idx, (ltext, lrect, _) in enumerate(self.labels):
                # Skip title labels (very wide, span full slide)
                lw = lrect[2] - lrect[0]
                if lw > 8.0:
                    continue  # title label

                ll, lt, lr, lb = lrect
                lcx = (ll + lr) / 2
                lcy = (lt + lb) / 2

                # Skip if arrow starts or ends near this label
                if (abs(bx - lcx) < (lr - ll) and abs(by - lcy) < (lb - lt)):
                    continue
                if (abs(ex - lcx) < (lr - ll) and abs(ey - lcy) < (lb - lt)):
                    continue

                if self._line_intersects_rect(bx, by, ex, ey, lrect, tol=0.03):
                    self.warnings.append(
                        f'Arrow ({bx:.2f},{by:.2f})→({ex:.2f},{ey:.2f}) '
                        f'passes through label "{ltext}"'
                    )

    def _line_intersects_rect(self, x1, y1, x2, y2, rect, tol=0.05):
        """Check if a line segment intersects a rectangle (shrunk by tol)."""
        rl, rt, rr, rb = rect
        # Shrink rect by tolerance
        rl += tol
        rt += tol
        rr -= tol
        rb -= tol
        if rr <= rl or rb <= rt:
            return False

        # Check if line segment passes through rectangle using
        # parametric clipping (Cohen-Sutherland-like)
        dx = x2 - x1
        dy = y2 - y1

        # For each edge, compute parameter t where line crosses
        t_min = 0.0
        t_max = 1.0

        for p, q in [(-dx, x1 - rl), (dx, rr - x1),
                      (-dy, y1 - rt), (dy, rb - y1)]:
            if abs(p) < 1e-10:
                if q < 0:
                    return False
            else:
                t = q / p
                if p < 0:
                    t_min = max(t_min, t)
                else:
                    t_max = min(t_max, t)

        return t_min <= t_max

    # ─── Check 5: Nested group containment ────────────────────────────────

    def _check_nested_group_containment(self):
        """Check that smaller groups are fully contained within larger groups."""
        if len(self.groups) < 2:
            return

        # Sort by area (largest first)
        sorted_groups = sorted(
            self.groups,
            key=lambda g: (g[1][2] - g[1][0]) * (g[1][3] - g[1][1]),
            reverse=True
        )

        for i in range(len(sorted_groups)):
            for j in range(i + 1, len(sorted_groups)):
                outer_name, outer_rect, _ = sorted_groups[i]
                inner_name, inner_rect, _ = sorted_groups[j]

                # Check if inner overlaps outer (they should either contain or not overlap)
                if rects_overlap(outer_rect, inner_rect, tol=0.0):
                    if not rect_contains(outer_rect, inner_rect, tol=0.3):
                        # Partially overlapping groups — ERROR
                        self.errors.append(
                            f'Nested group overlap: "{inner_name}" {rect_str(inner_rect)} '
                            f'partially overlaps "{outer_name}" {rect_str(outer_rect)} '
                            f'but is not fully contained. Child groups must be inside parent groups'
                        )

    # ─── Check 6: Boundary checks ────────────────────────────────────────

    def _check_boundary(self):
        for iname, irect, _ in self.icons:
            if irect[2] > SLIDE_W_IN + 0.05 or irect[3] > SLIDE_H_IN + 0.05:
                self.errors.append(
                    f'Icon {iname} extends beyond slide boundary at {rect_str(irect)}'
                )
            if irect[0] < -0.05 or irect[1] < -0.05:
                self.errors.append(
                    f'Icon {iname} extends beyond slide boundary (negative coords) at {rect_str(irect)}'
                )

        for gname, grect, _ in self.groups:
            if grect[2] > SLIDE_W_IN + 0.05:
                self.errors.append(
                    f'Group "{gname}" extends beyond right slide edge at x={grect[2]:.1f}"'
                )
            if grect[3] > SLIDE_H_IN + 0.05:
                self.errors.append(
                    f'Group "{gname}" extends beyond bottom slide edge at y={grect[3]:.1f}"'
                )

    # ─── Check 7: Font sizes ─────────────────────────────────────────────

    def _check_font(self, shape):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    pt = run.font.size / 12700
                    if pt < MIN_FONT_PT:
                        self.warnings.append(
                            f'Text "{run.text[:30]}" font size {pt:.0f}pt < minimum {MIN_FONT_PT}pt'
                        )
                    if pt > MAX_FONT_PT:
                        self.warnings.append(
                            f'Text "{run.text[:30]}" font size {pt:.0f}pt > maximum {MAX_FONT_PT}pt'
                        )

    # ─── Check 8: Group colors ───────────────────────────────────────────

    def _check_group_color(self, color_hex, rect):
        if not color_hex:
            return
        if color_hex in GROUP_COLORS:
            self.info.append(f'Group: {GROUP_COLORS[color_hex]} (#{color_hex}) at ({rect[0]:.1f}", {rect[1]:.1f}")')
        elif color_hex not in {'ffffff', '000000', '7d8998'}:
            self.warnings.append(
                f'Group border color #{color_hex} is not a standard AWS group color. '
                f'Standard: {", ".join("#" + c for c in sorted(GROUP_COLORS.keys()))}'
            )

    # ─── Check 9: Arrow styles ───────────────────────────────────────────

    def _check_arrow_styles(self):
        for ctype, shape in self.connectors:
            elem = shape._element
            ln = elem.find('.//' + qn('a:ln'))
            if ln is not None:
                prstDash = ln.find(qn('a:prstDash'))
                if prstDash is not None:
                    dash_val = prstDash.get('val', '')
                    if dash_val and dash_val != 'solid':
                        self.errors.append(
                            f'Connector uses dashed line style "{dash_val}". '
                            f'All arrows must be SOLID per AWS guidelines'
                        )

    # ─── Check 10: Group proportionality ──────────────────────────────────

    def _check_group_proportionality(self):
        """Check that groups have reasonable proportions."""
        for gname, grect, _ in self.groups:
            gl, gt, gr, gb = grect
            gw = gr - gl
            gh = gb - gt

            # Minimum group size: at least 0.5" in each dimension
            if gw < 0.5 or gh < 0.5:
                self.errors.append(
                    f'Group "{gname}" is too small ({gw:.2f}" x {gh:.2f}"). '
                    f'Minimum group size is 0.5" x 0.5"'
                )

            # Extreme aspect ratio: warn if width/height ratio > 6:1 or < 1:6
            if gw > 0 and gh > 0:
                ratio = max(gw / gh, gh / gw)
                if ratio > 6:
                    self.warnings.append(
                        f'Group "{gname}" has extreme aspect ratio {ratio:.1f}:1 '
                        f'({gw:.2f}" x {gh:.2f}"). Consider more balanced dimensions'
                    )

            # Group must be large enough to contain at least one icon
            # Service icon = 0.83*scale ≈ 0.7" + label ≈ 1.1" total height
            min_content_h = 0.8
            min_content_w = 0.8
            if gw < min_content_w or gh < min_content_h:
                self.warnings.append(
                    f'Group "{gname}" ({gw:.2f}" x {gh:.2f}") may be too small '
                    f'to contain icons comfortably'
                )

    # ─── Check 11: Group padding ─────────────────────────────────────────

    def _check_group_padding(self):
        """Check that icons have minimum padding from group borders."""
        MIN_PADDING = 0.15  # minimum 0.15" from group border

        for iname, irect, _ in self.icons:
            il, it, ir, ib = irect

            # Find the smallest group that contains this icon
            best_group = None
            best_area = float('inf')
            for gname, grect, _ in self.groups:
                gl, gt, gr, gb = grect
                # Check if icon center is inside group
                icx = (il + ir) / 2
                icy = (it + ib) / 2
                if gl <= icx <= gr and gt <= icy <= gb:
                    area = (gr - gl) * (gb - gt)
                    if area < best_area:
                        best_area = area
                        best_group = (gname, grect)

            if best_group:
                gname, (gl, gt, gr, gb) = best_group
                # Check padding from each border
                pad_left = il - gl
                pad_top = it - gt
                pad_right = gr - ir
                pad_bottom = gb - ib

                min_pad = min(pad_left, pad_top, pad_right, pad_bottom)
                if min_pad < MIN_PADDING:
                    side = 'left' if pad_left == min_pad else \
                           'top' if pad_top == min_pad else \
                           'right' if pad_right == min_pad else 'bottom'
                    self.warnings.append(
                        f'Icon {iname} is only {min_pad:.2f}" from {side} border '
                        f'of group "{gname}". Minimum padding: {MIN_PADDING}"'
                    )

    # ─── Check 12: Icon spacing ──────────────────────────────────────────

    def _check_icon_spacing(self):
        """Check minimum spacing between adjacent icons."""
        MIN_SPACING = 0.3  # minimum 0.3" between icon edges

        for i in range(len(self.icons)):
            for j in range(i + 1, len(self.icons)):
                n1, r1, _ = self.icons[i]
                n2, r2, _ = self.icons[j]

                # Calculate minimum edge-to-edge distance
                # Horizontal distance
                if r1[2] < r2[0]:
                    hdist = r2[0] - r1[2]
                elif r2[2] < r1[0]:
                    hdist = r1[0] - r2[2]
                else:
                    hdist = 0  # overlapping horizontally

                # Vertical distance
                if r1[3] < r2[1]:
                    vdist = r2[1] - r1[3]
                elif r2[3] < r1[1]:
                    vdist = r1[1] - r2[3]
                else:
                    vdist = 0  # overlapping vertically

                # If they overlap in one axis, check the other
                if hdist == 0 and vdist == 0:
                    continue  # already caught by overlap check
                elif hdist == 0:
                    dist = vdist
                elif vdist == 0:
                    dist = hdist
                else:
                    dist = (hdist ** 2 + vdist ** 2) ** 0.5

                if dist < MIN_SPACING and dist > 0:
                    self.warnings.append(
                        f'Icons {n1} and {n2} are only {dist:.2f}" apart. '
                        f'Minimum spacing: {MIN_SPACING}"'
                    )

    # ─── Report ───────────────────────────────────────────────────────────

    def report(self):
        print(f'\n{"="*60}')
        print(f'AWS Architecture Diagram Validation Report')
        print(f'File: {self.filepath}')
        print(f'{"="*60}')

        if self.info:
            print(f'\n--- INFO ---')
            for msg in self.info:
                print(f'  ℹ {msg}')

        if self.errors:
            print(f'\n--- ERRORS ({len(self.errors)}) ---')
            for msg in self.errors:
                print(f'  ✗ {msg}')

        if self.warnings:
            print(f'\n--- WARNINGS ({len(self.warnings)}) ---')
            for msg in self.warnings:
                print(f'  ⚠ {msg}')

        if not self.errors and not self.warnings:
            print(f'\n✓ All checks passed!')
        elif not self.errors:
            print(f'\n✓ Passed with {len(self.warnings)} warning(s)')
        else:
            print(f'\n✗ FAILED with {len(self.errors)} error(s), {len(self.warnings)} warning(s)')

        return len(self.errors) == 0


def main():
    parser = argparse.ArgumentParser(description='Validate AWS architecture PPTX diagram')
    parser.add_argument('file', help='Path to .pptx file')
    args = parser.parse_args()

    validator = DiagramValidator(args.file)
    validator.validate()
    passed = validator.report()
    sys.exit(0 if passed else 1)


if __name__ == '__main__':
    main()
